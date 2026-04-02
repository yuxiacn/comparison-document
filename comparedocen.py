#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Document Comparison Tool
Supports formats: PDF, DOCX, PPTX, TXT
Output: Comparison_File1_VS_File2.docx (landscape, diff-only view, word-level highlighting)
"""

import sys
import os
import re
import difflib
from datetime import datetime
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# 版本号
# Format: V{major}.{minor} Build{YYYYMMDD}.{revision}
# Update Rules:
#   - Major updates: increment major version (e.g., V2.0 → V3.0)
#   - Feature changes: increment minor version (e.g., V2.0 → V2.1)
#   - Bug fixes: increment revision (e.g., V2.0 Build20260403.1 → V2.0 Build20260403.2)
#   - Update date and revision with each modification
VERSION = "V2.0 Build20260403.1"

# 读取器注册表
READERS = {}


def register_reader(ext):
    def decorator(func):
        READERS[ext.lower()] = func
        return func
    return decorator


def estimate_visual_lines(text, chars_per_line=80):
    """
    估算文本的视觉行数（模拟Word自动换行）
    Consider mixed Chinese/English text, Chinese chars ~2x English width
    """
    if not text.strip():
        return 1  # 空行也算一行
    
    # 计算等效字符数（中文字符算2个宽度）
    effective_chars = 0
    for char in text:
        if '\u4e00' <= char <= '\u9fff':  # 中文字符
            effective_chars += 2
        else:
            effective_chars += 1
    
    # 计算需要的行数
    lines_needed = max(1, (effective_chars + chars_per_line - 1) // chars_per_line)
    return lines_needed


@register_reader('.txt')
def read_txt(path, merge_lines=False):
    """读取TXT，返回内容列表和位置信息(paragraphs号, page number=1, None)"""
    with open(path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    if merge_lines:
        # 合并连续非空行（处理自动换行的情况）
        lines = text.splitlines()
        paragraphs = []
        current_para = []
        
        for line in lines:
            stripped = line.rstrip()
            if not stripped:
                if current_para:
                    paragraphs.append(' '.join(current_para))
                    current_para = []
            else:
                current_para.append(stripped)
        
        if current_para:
            paragraphs.append(' '.join(current_para))
        
        lines = paragraphs
    else:
        lines = [ln.rstrip() for ln in text.splitlines() if ln.strip()]
    
    # TXT没有真实page number和line number
    location_info = [(i + 1, 1, None) for i in range(len(lines))]
    return lines, location_info


def estimate_paragraph_pages(doc):
    """
    估算每个paragraphs的page number
    Returns: [page_number, ...] 与paragraphs一一对应
    """
    # 获取页面设置
    section = doc.sections[0] if doc.sections else None
    
    if section:
        # 页面高度（英寸）
        page_height = section.page_height.inches if section.page_height else 11
        # 上下边距
        top_margin = section.top_margin.inches if section.top_margin else 1
        bottom_margin = section.bottom_margin.inches if section.bottom_margin else 1
        # 可用高度
        available_height = page_height - top_margin - bottom_margin
    else:
        available_height = 9  # 默认可用高度
    
    page_numbers = []
    current_page = 1
    current_page_used_height = 0
    
    # 估算每英寸高度可容纳的 12pt 文字行数（约5-6行）
    lines_per_inch = 5.5
    
    for para in doc.paragraphs:
        text = para.text.rstrip()
        
        # 获取字体大小
        font_size = 12
        try:
            if para.runs and para.runs[0].font.size:
                font_size = para.runs[0].font.size.pt
        except:
            pass
        
        # 计算该paragraphs占用的高度（英寸）
        # 字体越大，行高越大
        line_height = (font_size / 12) * (1 / lines_per_inch)
        
        # 估算paragraphs行数（简单估算）
        effective_chars = sum(2 if '\u4e00' <= c <= '\u9fff' else 1 for c in text)
        chars_per_line = 80  # 简化估算
        lines_needed = max(1, (effective_chars + chars_per_line - 1) // chars_per_line)
        
        para_height = lines_needed * line_height
        
        # paragraphs前后间距
        space_before = 0
        space_after = 0
        try:
            if para.paragraph_format.space_before:
                space_before = para.paragraph_format.space_before.pt / 72  # 转换为英寸
            if para.paragraph_format.space_after:
                space_after = para.paragraph_format.space_after.pt / 72
        except:
            pass
        
        total_height = space_before + para_height + space_after
        
        # 检查是否跨页
        if current_page_used_height + total_height > available_height:
            current_page += 1
            current_page_used_height = total_height
        else:
            current_page_used_height += total_height
        
        page_numbers.append(current_page)
    
    return page_numbers


@register_reader('.docx')
def read_docx(path, use_precise=True):
    """
    读取docx，返回内容列表和位置信息
    Location format: [(paragraph_number, page_number, line_number), ...]
    DOCX has no visual line numbers, so line_number is None
    """
    doc = Document(path)
    lines = [para.text.rstrip() for para in doc.paragraphs]
    
    # 估算page number
    page_numbers = estimate_paragraph_pages(doc)
    
    # 构建位置信息 (paragraphs号, page number, None)
    location_info = []
    for i, page_num in enumerate(page_numbers):
        location_info.append((i + 1, page_num, None))
    
    return lines, location_info


@register_reader('.pdf')
def read_pdf(path, merge_lines=True, merge_across_pages=True):
    """
    读取PDF，返回内容列表和位置信息(paragraphs号, page number, line number)
    
    Args:
        merge_lines: whether to merge consecutive non-empty lines into a paragraph
        merge_across_pages: whether to merge paragraphs across pages
    
    Returns:
        paragraphs: list of text content
        location_info: [(paragraph_num, page_num, line_num), ...]
    """
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("请安装 pdfplumber 以支持 PDF 读取: pip install pdfplumber")
    
    paragraphs = []
    location_info = []
    paragraph_counter = 0
    
    # 用于识别page number的模式：单独的数字行（1-4位数字）
    import re
    page_number_pattern = re.compile(r'^\s*\d{1,4}\s*$')
    # 用于识别行首的视觉line number：行开头的数字（空格+数字+空格或句点）
    line_number_pattern = re.compile(r'^(\s*\d+)[\.\s]\s*')
    
    # 先收集所有页面的原始行信息
    all_pages_lines = []  # [(page_num, lines), ...]
    
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if not text:
                all_pages_lines.append((page_num, []))
                continue
            
            page_lines = text.splitlines()
            processed_lines = []
            
            for line in page_lines:
                line = line.rstrip()
                if not line:
                    continue
                # 过滤page number行
                if page_number_pattern.match(line):
                    continue
                
                # 提取行首的视觉line number
                visual_line_num = None
                match = line_number_pattern.match(line)
                if match:
                    try:
                        visual_line_num = int(match.group(1).strip())
                        line = line[match.end():].lstrip()
                    except ValueError:
                        pass
                
                if line:
                    processed_lines.append((line, visual_line_num))
            
            all_pages_lines.append((page_num, processed_lines))
    
    # 跨页paragraphs合并处理
    if merge_across_pages and merge_lines:
        # 将所有页面的行合并，然后统一处理paragraphs
        all_lines = []
        for page_num, lines in all_pages_lines:
            for line, line_num in lines:
                all_lines.append((line, line_num, page_num))
        
        # 统一处理paragraphs合并
        current_paragraph = []
        current_line_num = None
        current_page = None
        
        i = 0
        while i < len(all_lines):
            line, visual_line_num, page_num = all_lines[i]
            
            # 检查是否是新paragraphs的开始
            is_new_para = False
            
            if not current_paragraph:
                # 第一个paragraphs开始
                is_new_para = True
            else:
                # 检查当前行是否属于新paragraphs
                stripped = line.lstrip()
                indent = len(line) - len(stripped)
                
                # 条件1: 明显缩进（4空格以上）
                if indent >= 4:
                    is_new_para = True
                # 条件2: 上一行以句号等结束，且当前行以大写或数字开头
                else:
                    prev_line = current_paragraph[-1]
                    if prev_line and prev_line[-1] in '.。!?':
                        if stripped and (stripped[0].isupper() or stripped[0].isdigit()):
                            is_new_para = True
                    # 条件3: 当前行是空行分隔（已过滤）或line number重置（明显变小）
                    if visual_line_num and current_line_num:
                        if visual_line_num < current_line_num and visual_line_num < 10:
                            is_new_para = True
            
            if is_new_para and current_paragraph:
                # 保存当前paragraphs
                paragraph_counter += 1
                para_text = ' '.join(current_paragraph)
                paragraphs.append(para_text)
                location_info.append((paragraph_counter, current_page or page_num, current_line_num))
                
                # 开始新paragraphs
                current_paragraph = [line]
                current_line_num = visual_line_num
                current_page = page_num
            else:
                # 继续当前paragraphs
                if not current_paragraph:
                    current_line_num = visual_line_num
                    current_page = page_num
                current_paragraph.append(line)
            
            i += 1
        
        # 保存最后一个paragraphs
        if current_paragraph:
            paragraph_counter += 1
            para_text = ' '.join(current_paragraph)
            paragraphs.append(para_text)
            location_info.append((paragraph_counter, current_page, current_line_num))
    
    else:
        # 逐页处理，不跨页合并
        for page_num, lines in all_pages_lines:
            if merge_lines:
                current_paragraph = []
                current_line_num = None
                
                for line, visual_line_num in lines:
                    if not line:
                        if current_paragraph:
                            paragraph_counter += 1
                            para_text = ' '.join(current_paragraph)
                            paragraphs.append(para_text)
                            location_info.append((paragraph_counter, page_num, current_line_num))
                            current_paragraph = []
                            current_line_num = None
                        continue
                    
                    # 检查新paragraphs
                    is_new_para = False
                    if current_paragraph:
                        stripped = line.lstrip()
                        indent = len(line) - len(stripped)
                        if indent >= 4:
                            is_new_para = True
                        prev_line = current_paragraph[-1]
                        if prev_line and prev_line[-1] in '.。!?':
                            if stripped and not stripped[0].islower():
                                is_new_para = True
                    
                    if is_new_para and current_paragraph:
                        paragraph_counter += 1
                        para_text = ' '.join(current_paragraph)
                        paragraphs.append(para_text)
                        location_info.append((paragraph_counter, page_num, current_line_num))
                        current_paragraph = [line]
                        current_line_num = visual_line_num
                    else:
                        if not current_paragraph:
                            current_line_num = visual_line_num
                        current_paragraph.append(line)
                
                # 保存最后一个paragraphs
                if current_paragraph:
                    paragraph_counter += 1
                    para_text = ' '.join(current_paragraph)
                    paragraphs.append(para_text)
                    location_info.append((paragraph_counter, page_num, current_line_num))
            else:
                # 不合并，每行独立
                for line, visual_line_num in lines:
                    if line:
                        paragraph_counter += 1
                        paragraphs.append(line)
                        location_info.append((paragraph_counter, page_num, visual_line_num))
    
    return paragraphs, location_info


@register_reader('.pptx')
def read_pptx(path):
    """读取PPTX，返回内容列表和位置信息(paragraphs号, 幻灯片号, None)"""
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    location_info = []
    paragraph_counter = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for ln in shape.text.splitlines():
                    line = ln.rstrip()
                    if line:  # 只处理非空行
                        paragraph_counter += 1
                        lines.append(line)
                        location_info.append((paragraph_counter, slide_num, None))
    
    return lines, location_info


def read_document(path, use_precise=True, merge_lines=True, merge_across_pages=True):
    """
    读取文档内容
    use_precise: for DOCX, whether to use precise visual line number calculation
    merge_lines: for PDF/TXT, whether to merge consecutive lines
    merge_across_pages: for PDF, whether to merge paragraphs across pages
    """
    ext = Path(path).suffix.lower()
    if ext not in READERS:
        raise ValueError(f"不支持的文件格式: {ext}。当前支持: {', '.join(READERS.keys())}")
    
    reader = READERS[ext]
    # 对不同格式传递相应参数
    if ext == '.docx':
        return reader(path, use_precise=use_precise)
    elif ext == '.pdf':
        return reader(path, merge_lines=merge_lines, merge_across_pages=merge_across_pages)
    elif ext == '.txt':
        return reader(path, merge_lines=merge_lines)
    else:
        return reader(path)


def set_run_font(run, font_name='Times New Roman', east_asia='宋体', size=Pt(10), bold=False):
    """统一设置中英文字体"""
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn('w:eastAsia'), east_asia)
    run.font.size = size
    if bold:
        run.font.bold = True


def add_colored_run(paragraph, text, rgb, bold=False):
    run = paragraph.add_run(text)
    set_run_font(run, bold=bold)
    run.font.color.rgb = RGBColor(*rgb)
    return run


def add_page_number(paragraph):
    """在paragraphs中添加page number字段"""
    run = paragraph.add_run()
    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')

    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = "PAGE"

    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'end')

    run._r.append(fldChar1)
    run._r.append(instrText)
    run._r.append(fldChar2)
    return run


def tokenize_text(text):
    """
    将文本拆分为单词和分隔符（空格/标点等）。
    Preserve all content for later reassembly.
    """
    return re.findall(r'\S+|\s+', text)


def split_sentences(text):
    """
    将文本按句子分割，保留分隔符。
    Sentence delimiters: Chinese period, English period, question mark, exclamation mark
    """
    if not text.strip():
        return []
    # 按句子分隔符分割，保留分隔符
    parts = re.split(r'([。！？.!?])', text)
    sentences = []
    current = ''
    for part in parts:
        current += part
        if re.match(r'[。！？.!?]$', part):
            # 遇到句子结束符，完成一个句子
            sentences.append(current)
            current = ''
    if current.strip():
        # 剩余内容（可能是不完整的句子）也作为一个句子
        sentences.append(current)
    return sentences


def word_diff_runs(text1, text2):
    """
    单词级精细差异分析，支持句子级缺失检测。
    Returns left_runs and right_runs, each is a list of (text, is_diff, is_placeholder).
    is_placeholder=True indicates a missing content placeholder [Missing Sentence]
    
    Logic:
    1. First split text into sentences
    2. Compare at sentence level
    3. When one side has a complete sentence (with delimiter) and the other does not, show placeholder
    """
    # 首先尝试句子级对比
    sents1 = split_sentences(text1)
    sents2 = split_sentences(text2)
    
    # 如果都能分割成句子，进行句子级对比
    if len(sents1) > 0 and len(sents2) > 0:
        sm_sents = difflib.SequenceMatcher(None, sents1, sents2, autojunk=False)
        opcodes = sm_sents.get_opcodes()
        
        # 检查是否有句子级别的 delete/insert
        has_sentence_level_diff = any(
            tag in ('delete', 'insert') for tag, _, _, _, _ in opcodes
        )
        
        if has_sentence_level_diff:
            left_runs = []
            right_runs = []
            
            for tag, i1, i2, j1, j2 in opcodes:
                if tag == 'equal':
                    for k in range(i1, min(i2, len(sents1))):
                        left_runs.append((sents1[k], False, False))
                    for k in range(j1, min(j2, len(sents2))):
                        right_runs.append((sents2[k], False, False))
                elif tag == 'replace':
                    # 替换：两边内容不同
                    max_len = max(i2 - i1, j2 - j1)
                    for k in range(max_len):
                        if i1 + k < i2 and i1 + k < len(sents1):
                            left_runs.append((sents1[i1 + k], True, False))
                        if j1 + k < j2 and j1 + k < len(sents2):
                            right_runs.append((sents2[j1 + k], True, False))
                elif tag == 'delete':
                    # 左边有句子，右边缺失整句
                    for k in range(i1, min(i2, len(sents1))):
                        left_runs.append((sents1[k], True, False))
                        right_runs.append(('[Missing Sentence]', True, True))
                elif tag == 'insert':
                    # 右边有句子，左边缺失整句
                    for k in range(j1, min(j2, len(sents2))):
                        left_runs.append(('[Missing Sentence]', True, True))
                        right_runs.append((sents2[k], True, False))
            
            # 合并连续的 [Missing Sentence] 占位符
            def merge_consecutive_placeholders(runs):
                if not runs:
                    return runs
                merged = [runs[0]]
                for i in range(1, len(runs)):
                    prev_text, prev_diff, prev_placeholder = merged[-1]
                    curr_text, curr_diff, curr_placeholder = runs[i]
                    # 如果前一个是占位符且当前也是占位符，跳过当前（合并）
                    if prev_placeholder and curr_placeholder:
                        continue
                    merged.append((curr_text, curr_diff, curr_placeholder))
                return merged

            left_runs = merge_consecutive_placeholders(left_runs)
            right_runs = merge_consecutive_placeholders(right_runs)
            
            return left_runs, right_runs
    
    # 回退到单词级对比（不显示占位符）
    tokens1 = tokenize_text(text1)
    tokens2 = tokenize_text(text2)
    sm = difflib.SequenceMatcher(None, tokens1, tokens2, autojunk=False)

    left_runs = []
    right_runs = []

    for tag, i1, i2, j1, j2 in sm.get_opcodes():
        if tag == 'equal':
            seg1 = ''.join(tokens1[i1:i2])
            left_runs.append((seg1, False, False))
            seg2 = ''.join(tokens2[j1:j2])
            right_runs.append((seg2, False, False))
        elif tag == 'replace':
            seg1 = ''.join(tokens1[i1:i2])
            left_runs.append((seg1, True, False))
            seg2 = ''.join(tokens2[j1:j2])
            right_runs.append((seg2, True, False))
        elif tag == 'delete':
            seg1 = ''.join(tokens1[i1:i2])
            left_runs.append((seg1, True, False))
        elif tag == 'insert':
            seg2 = ''.join(tokens2[j1:j2])
            right_runs.append((seg2, True, False))

    if not left_runs:
        left_runs.append(('', False, False))
    if not right_runs:
        right_runs.append(('', False, False))

    return left_runs, right_runs


def get_word_line_number_offset(doc_path):
    """
    尝试读取 Word 文档的line number设置，返回line number起始偏移量。
    如果文档启用了line number，返回line number起始值；否则返回 None。
    """
    try:
        doc = Document(doc_path)
        for section in doc.sections:
            # 尝试获取line number设置
            sectPr = section._sectPr
            if hasattr(sectPr, 'lnNumType') and sectPr.lnNumType is not None:
                # 文档启用了line number
                lnNumType = sectPr.lnNumType
                start = getattr(lnNumType, 'start', 1)
                return int(start) if start else 1
    except Exception:
        pass
    return None


def build_diff_report(lines1, lines2, location_info1=None, location_info2=None):
    """
    使用 difflib.SequenceMatcher 分析差异，返回仅包含差异行的数据。
    Each element is (tag, left_loc, left_text, right_loc, right_text)
    tag values: replace, delete, insert
    location 格式: (paragraph_number, page_number) 或 None
    """
    sm = difflib.SequenceMatcher(None, lines1, lines2, autojunk=False)
    opcodes = sm.get_opcodes()
    
    # 如果没有提供位置信息，使用默认的paragraphs索引+1
    if location_info1 is None:
        location_info1 = [(i + 1, 1) for i in range(len(lines1))]
    if location_info2 is None:
        location_info2 = [(i + 1, 1) for i in range(len(lines2))]

    rows = []
    for tag, i1, i2, j1, j2 in opcodes:
        if tag == 'equal':
            continue
        elif tag == 'replace':
            max_len = max(i2 - i1, j2 - j1)
            for k in range(max_len):
                left_exists = i1 + k < i2
                right_exists = j1 + k < j2
                left_loc = location_info1[i1 + k] if left_exists else None
                right_loc = location_info2[j1 + k] if right_exists else None
                ltext = lines1[i1 + k] if left_exists else ""
                rtext = lines2[j1 + k] if right_exists else ""
                if left_exists and right_exists:
                    rows.append(('replace', left_loc, ltext, right_loc, rtext))
                elif left_exists and not right_exists:
                    rows.append(('delete', left_loc, ltext, None, ""))
                elif not left_exists and right_exists:
                    rows.append(('insert', None, "", right_loc, rtext))
        elif tag == 'delete':
            for k in range(i2 - i1):
                left_loc = location_info1[i1 + k]
                ltext = lines1[i1 + k]
                rows.append(('delete', left_loc, ltext, None, ""))
        elif tag == 'insert':
            for k in range(j2 - j1):
                right_loc = location_info2[j1 + k]
                rtext = lines2[j1 + k]
                rows.append(('insert', None, "", right_loc, rtext))

    return rows


def set_cell_width(cell, width_inches):
    """通过底层 XML 强制设置单元格宽度"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcW = OxmlElement('w:tcW')
    tcW.set(qn('w:w'), str(int(width_inches * 1440)))
    tcW.set(qn('w:type'), 'dxa')
    tcPr.append(tcW)


def set_table_column_widths(table, widths_inches):
    """
    通过底层 XML 精确设置表格各列宽度及总宽度。
    widths_inches: list of width values in inches for each column
    """
    tbl = table._tbl
    tblGrid = tbl.tblGrid
    gridCols = tblGrid.gridCol_lst
    for idx, w in enumerate(widths_inches):
        if idx < len(gridCols):
            gridCols[idx].set(qn('w:w'), str(int(w * 1440)))
            gridCols[idx].set(qn('w:type'), 'dxa')
        if idx < len(table.columns):
            table.columns[idx].width = Inches(w)

    tblPr = tbl.tblPr
    tblW_list = tblPr.xpath('./w:tblW')
    if tblW_list:
        tblW = tblW_list[0]
    else:
        tblW = OxmlElement('w:tblW')
        tblPr.append(tblW)
    total_width = sum(widths_inches)
    tblW.set(qn('w:w'), str(int(total_width * 1440)))
    tblW.set(qn('w:type'), 'dxa')


def generate_docx(rows, name1, name2, output_path):
    doc = Document()

    # 设置为横向（Landscape），并应用标准 1 英寸页边距
    section = doc.sections[0]
    section.orientation = WD_ORIENT.LANDSCAPE
    new_width, new_height = section.page_height, section.page_width
    section.page_width = new_width
    section.page_height = new_height
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)

    # 添加页脚page number
    footer = section.footer
    footer_para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = add_page_number(footer_para)
    set_run_font(run, size=Pt(9))

    # 颜色定义
    color_left = (0, 112, 192)    # 蓝色
    color_right = (255, 0, 0)     # 红色
    color_mark_replace = (255, 140, 0)  # 橙色
    color_gray = (100, 100, 100)
    color_placeholder = (0, 176, 80)  # 绿色，用于缺失句子

    # First line: Document Comparison Report
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_title.add_run('Document Comparison Report')
    set_run_font(run, font_name='Noto Serif', east_asia='Noto Serif', size=Pt(18), bold=True)

    # 第二行：两个文档名称
    p_names = doc.add_paragraph()
    p_names.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_names.add_run(name1)
    set_run_font(run, size=Pt(14), bold=True)
    run.font.color.rgb = RGBColor(*color_left)
    run = p_names.add_run(' VS ')
    set_run_font(run, size=Pt(14), bold=True)
    run = p_names.add_run(name2)
    set_run_font(run, size=Pt(14), bold=True)
    run.font.color.rgb = RGBColor(*color_right)

    # 第三行：生成日期时间
    now_str = datetime.now().strftime('%Y-%m-%d %H:%M')
    p_time = doc.add_paragraph()
    p_time.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_time.add_run(now_str)
    set_run_font(run, size=Pt(12), bold=True)
    run.font.color.rgb = RGBColor(*color_gray)

    # 时间行与图例行之间空一行
    doc.add_paragraph()

    # 图例说明（宋体）
    p = doc.add_paragraph()
    run = p.add_run('Legend: ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run(' = : identical content (hidden)  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('≠ ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_mark_replace)
    run = p.add_run(': modified content  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('- ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_left)
    run = p.add_run(': deleted content  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('+ ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_right)
    run = p.add_run(': added content  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('P: page number  L: line number')
    set_run_font(run, east_asia='宋体', bold=False)

    if not rows:
        p = doc.add_paragraph()
        run = p.add_run('Both documents are identical, no differences found.')
        set_run_font(run, size=Pt(12))
        run.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.save(output_path)
        return

    # 表格：5列（总宽度 8.8 英寸，适应 1 英寸边距的横向页面）
    table = doc.add_table(rows=1, cols=5)
    table.style = 'Table Grid'
    table.autofit = False
    table.allow_autofit = False

    col_widths = [0.90, 3.35, 0.70, 0.90, 3.35]
    set_table_column_widths(table, col_widths)

    # 设置表头 - 改为显示page number-paragraphs号
    hdr_cells = table.rows[0].cells
    headers = ['Location', name1, 'Mark', 'Location', name2]
    header_colors = [None, color_left, None, None, color_right]
    for idx, text in enumerate(headers):
        cell = hdr_cells[idx]
        set_cell_width(cell, col_widths[idx])
        p = cell.paragraphs[0]
        p.clear()
        run = p.add_run(text)
        set_run_font(run, bold=True)
        if header_colors[idx]:
            run.font.color.rgb = RGBColor(*header_colors[idx])
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for tag, left_loc, ltext, right_loc, rtext in rows:
        row_cells = table.add_row().cells

        # 位置1 - 格式: Ppage number-Lline number（或 Ppage number-paragraphs号）
        cell = row_cells[0]
        set_cell_width(cell, col_widths[0])
        p = cell.paragraphs[0]
        p.clear()
        if left_loc is not None:
            # 处理三元组 (para_num, page_num, line_num) 或二元组 (para_num, page_num)
            if len(left_loc) >= 3:
                para_num, page_num, line_num = left_loc[0], left_loc[1], left_loc[2]
                if line_num:
                    loc_text = f"P{page_num}-L{line_num}"
                else:
                    loc_text = f"P{page_num}-{para_num}"
            else:
                para_num, page_num = left_loc[0], left_loc[1]
                loc_text = f"P{page_num}-{para_num}"
        else:
            loc_text = ""
        run = p.add_run(loc_text)
        set_run_font(run, size=Pt(8))
        run.font.color.rgb = RGBColor(*color_gray)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 位置2 - 格式: Ppage number-Lline number（或 Ppage number-paragraphs号）
        cell = row_cells[3]
        set_cell_width(cell, col_widths[3])
        p = cell.paragraphs[0]
        p.clear()
        if right_loc is not None:
            # 处理三元组或二元组
            if len(right_loc) >= 3:
                para_num, page_num, line_num = right_loc[0], right_loc[1], right_loc[2]
                if line_num:
                    loc_text = f"P{page_num}-L{line_num}"
                else:
                    loc_text = f"P{page_num}-{para_num}"
            else:
                para_num, page_num = right_loc[0], right_loc[1]
                loc_text = f"P{page_num}-{para_num}"
        else:
            loc_text = ""
        run = p.add_run(loc_text)
        set_run_font(run, size=Pt(8))
        run.font.color.rgb = RGBColor(*color_gray)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 内容1
        cell = row_cells[1]
        set_cell_width(cell, col_widths[1])
        p1 = cell.paragraphs[0]
        p1.clear()
        if tag == 'replace':
            runs1, _ = word_diff_runs(ltext, rtext)
            for text_seg, is_diff, is_placeholder in runs1:
                if is_placeholder:
                    # 缺失占位符：绿色粗体
                    add_colored_run(p1, text_seg, color_placeholder, bold=True)
                elif is_diff:
                    add_colored_run(p1, text_seg, color_left, bold=True)
                else:
                    run = p1.add_run(text_seg)
                    set_run_font(run)
        elif tag == 'delete':
            add_colored_run(p1, ltext, color_left, bold=True)
            # 右侧整行空白，不添加任何标识
        elif tag == 'insert':
            # 左侧整行空白，不添加任何标识
            pass
        else:
            run = p1.add_run(ltext)
            set_run_font(run)

        # 标记
        cell = row_cells[2]
        set_cell_width(cell, col_widths[2])
        p_mark = cell.paragraphs[0]
        p_mark.clear()
        p_mark.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if tag == 'replace':
            add_colored_run(p_mark, '≠', color_mark_replace, bold=True)
        elif tag == 'delete':
            add_colored_run(p_mark, '-', color_left, bold=True)
        elif tag == 'insert':
            add_colored_run(p_mark, '+', color_right, bold=True)

        # 内容2
        cell = row_cells[4]
        set_cell_width(cell, col_widths[4])
        p2 = cell.paragraphs[0]
        p2.clear()
        if tag == 'replace':
            _, runs2 = word_diff_runs(ltext, rtext)
            for text_seg, is_diff, is_placeholder in runs2:
                if is_placeholder:
                    # 缺失占位符：绿色粗体
                    add_colored_run(p2, text_seg, color_placeholder, bold=True)
                elif is_diff:
                    add_colored_run(p2, text_seg, color_right, bold=True)
                else:
                    run = p2.add_run(text_seg)
                    set_run_font(run)
        elif tag == 'insert':
            add_colored_run(p2, rtext, color_right, bold=True)
        elif tag == 'delete':
            # 删除行：右侧整行空白，不添加任何标识
            pass
        else:
            run = p2.add_run(rtext)
            set_run_font(run)

    doc.save(output_path)


def main():
    import argparse
    
    # 程序开头说明
    print("=" * 70)
    print(f"compare_docs.py - Document Difference Comparison Tool {VERSION}")
    print("Supports PDF/DOCX/PPTX/TXT comparison, generates Word diff report")
    print("Author: Yu Xia  E-mail: yuxiacn@qq.com")
    print("=" * 70)
    print()
    
    parser = argparse.ArgumentParser(
        description='Document Comparison Tool - Compare two documents and generate Word diff report',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
支持格式:
  PDF   - 支持跨页paragraphs合并、page number过滤、视觉line number提取
  DOCX  - 支持paragraphs级对比、估算page number
  PPTX  - 幻灯片文本对比
  TXT   - 纯文本对比

示例:
  python compare_docs.py paper1.pdf paper2.pdf
  python compare_docs.py report1.docx report2.docx --calibrate
  python compare_docs.py doc1.txt doc2.txt --no-merge
        '''
    )
    parser.add_argument('file1', help='First document path')
    parser.add_argument('file2', help='Second document path')
    parser.add_argument('--calibrate', action='store_true', help='Calibration mode: output paragraph location info for debugging')
    parser.add_argument('--no-merge', action='store_true', help='PDF/TXT files: do not merge consecutive lines (compare by original lines)')
    parser.add_argument('--no-page-merge', action='store_true', help='PDF: do not merge paragraphs across pages (process each page independently)')
    
    args = parser.parse_args()
    
    file1 = args.file1
    file2 = args.file2
    use_precise = True
    merge_lines = not args.no_merge  # 默认合并，--no-merge 时关闭
    merge_across_pages = not args.no_page_merge  # 默认跨页合并，--no-page-merge 时关闭
    
    if not os.path.exists(file1):
        print(f"Error: File does not exist: {file1}")
        sys.exit(1)
    if not os.path.exists(file2):
        print(f"Error: File does not exist: {file2}")
        sys.exit(1)

    name1 = Path(file1).stem
    name2 = Path(file2).stem
    output_name = f"Comparison_{name1}_VS_{name2}.docx"
    output_path = os.path.join(os.getcwd(), output_name)

    print(f"Reading {file1} ...")
    result1 = read_document(file1, use_precise=use_precise, merge_lines=merge_lines, merge_across_pages=merge_across_pages)
    if isinstance(result1, tuple):
        lines1, location_info1 = result1
    else:
        lines1 = result1
        location_info1 = None
    print(f"  Total {len(lines1)} paragraphs")
    
    # 校准模式：显示前几个paragraphs的位置信息
    if args.calibrate and location_info1:
        print("\n校准信息（前5个paragraphs）:")
        for i in range(min(5, len(lines1))):
            loc = location_info1[i]
            para_num, page_num, line_num = loc[0], loc[1], loc[2] if len(loc) > 2 else None
            if line_num:
                print(f"  paragraphs {i+1}: Page{page_num}-L{line_num}")
            else:
                print(f"  paragraphs {i+1}: Page{page_num}-{para_num}")
        print()

    print(f"Reading {file2} ...")
    result2 = read_document(file2, use_precise=use_precise, merge_lines=merge_lines, merge_across_pages=merge_across_pages)
    if isinstance(result2, tuple):
        lines2, location_info2 = result2
    else:
        lines2 = result2
        location_info2 = None
    print(f"  Total {len(lines2)} paragraphs")

    print("Analyzing differences ...")
    rows = build_diff_report(lines1, lines2, location_info1, location_info2)

    print("Generating report ...")
    generate_docx(rows, name1, name2, output_path)

    diff_count = len(rows)
    print(f"Diff row count: {diff_count}")
    print(f"Comparison report saved to: {output_path}")


if __name__ == '__main__':
    main()
