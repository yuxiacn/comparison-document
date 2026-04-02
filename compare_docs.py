#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档对比工具
支持格式：docx, pdf, pptx, txt
输出：Comparison_文件名1_VS_文件名2.docx（横向页面，仅显示差异行，单词级精确高亮）
"""

import sys
import os
import re
import difflib
from datetime import datetime
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# 读取器注册表
READERS = {}


def register_reader(ext):
    def decorator(func):
        READERS[ext.lower()] = func
        return func
    return decorator


def estimate_visual_lines(text, chars_per_line=80):
    """
    估算文本的视觉行数（模拟Word自动换行）
    考虑中英文混排，中文字符占约2个英文字符宽度
    """
    if not text.strip():
        return 1  # 空行也算一行
    
    # 计算等效字符数（中文字符算2个宽度）
    effective_chars = 0
    for char in text:
        if '\u4e00' <= char <= '\u9fff':  # 中文字符
            effective_chars += 2
        else:
            effective_chars += 1
    
    # 计算需要的行数
    lines_needed = max(1, (effective_chars + chars_per_line - 1) // chars_per_line)
    return lines_needed


@register_reader('.txt')
def read_txt(path, merge_lines=False):
    """读取TXT，返回内容列表和位置信息(段落号, 页码=1, None)"""
    with open(path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    if merge_lines:
        # 合并连续非空行（处理自动换行的情况）
        lines = text.splitlines()
        paragraphs = []
        current_para = []
        
        for line in lines:
            stripped = line.rstrip()
            if not stripped:
                if current_para:
                    paragraphs.append(' '.join(current_para))
                    current_para = []
            else:
                current_para.append(stripped)
        
        if current_para:
            paragraphs.append(' '.join(current_para))
        
        lines = paragraphs
    else:
        lines = [ln.rstrip() for ln in text.splitlines() if ln.strip()]
    
    # TXT没有真实页码和行号
    location_info = [(i + 1, 1, None) for i in range(len(lines))]
    return lines, location_info


def estimate_paragraph_pages(doc):
    """
    估算每个段落的页码
    返回: [page_number, ...] 与段落一一对应
    """
    # 获取页面设置
    section = doc.sections[0] if doc.sections else None
    
    if section:
        # 页面高度（英寸）
        page_height = section.page_height.inches if section.page_height else 11
        # 上下边距
        top_margin = section.top_margin.inches if section.top_margin else 1
        bottom_margin = section.bottom_margin.inches if section.bottom_margin else 1
        # 可用高度
        available_height = page_height - top_margin - bottom_margin
    else:
        available_height = 9  # 默认可用高度
    
    page_numbers = []
    current_page = 1
    current_page_used_height = 0
    
    # 估算每英寸高度可容纳的 12pt 文字行数（约5-6行）
    lines_per_inch = 5.5
    
    for para in doc.paragraphs:
        text = para.text.rstrip()
        
        # 获取字体大小
        font_size = 12
        try:
            if para.runs and para.runs[0].font.size:
                font_size = para.runs[0].font.size.pt
        except:
            pass
        
        # 计算该段落占用的高度（英寸）
        # 字体越大，行高越大
        line_height = (font_size / 12) * (1 / lines_per_inch)
        
        # 估算段落行数（简单估算）
        effective_chars = sum(2 if '\u4e00' <= c <= '\u9fff' else 1 for c in text)
        chars_per_line = 80  # 简化估算
        lines_needed = max(1, (effective_chars + chars_per_line - 1) // chars_per_line)
        
        para_height = lines_needed * line_height
        
        # 段落前后间距
        space_before = 0
        space_after = 0
        try:
            if para.paragraph_format.space_before:
                space_before = para.paragraph_format.space_before.pt / 72  # 转换为英寸
            if para.paragraph_format.space_after:
                space_after = para.paragraph_format.space_after.pt / 72
        except:
            pass
        
        total_height = space_before + para_height + space_after
        
        # 检查是否跨页
        if current_page_used_height + total_height > available_height:
            current_page += 1
            current_page_used_height = total_height
        else:
            current_page_used_height += total_height
        
        page_numbers.append(current_page)
    
    return page_numbers


@register_reader('.docx')
def read_docx(path, use_precise=True):
    """
    读取docx，返回内容列表和位置信息
    位置信息格式: [(paragraph_number, page_number, line_number), ...]
    DOCX 没有视觉行号，所以 line_number 为 None
    """
    doc = Document(path)
    lines = [para.text.rstrip() for para in doc.paragraphs]
    
    # 估算页码
    page_numbers = estimate_paragraph_pages(doc)
    
    # 构建位置信息 (段落号, 页码, None)
    location_info = []
    for i, page_num in enumerate(page_numbers):
        location_info.append((i + 1, page_num, None))
    
    return lines, location_info


@register_reader('.pdf')
def read_pdf(path, merge_lines=True):
    """
    读取PDF，返回内容列表和位置信息(段落号, 页码, 行号)
    
    参数:
        merge_lines: 是否将连续的非空行合并为一个段落
    
    返回:
        paragraphs: 文本内容列表
        location_info: [(paragraph_num, page_num, line_num), ...]
    """
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("请安装 pdfplumber 以支持 PDF 读取: pip install pdfplumber")
    
    paragraphs = []
    location_info = []
    paragraph_counter = 0
    
    # 用于识别页码的模式：单独的数字行（1-4位数字）
    import re
    page_number_pattern = re.compile(r'^\s*\d{1,4}\s*$')
    # 用于识别行首的视觉行号：行开头的数字（空格+数字+空格或句点）
    line_number_pattern = re.compile(r'^(\s*\d+)[\.\s]\s*')
    
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if not text:
                continue
                
            page_lines = text.splitlines()
            
            if merge_lines:
                # 智能段落合并
                current_paragraph = []
                current_line_num = None  # 记录段落的起始行号
                
                for i, line in enumerate(page_lines):
                    line = line.rstrip()
                    if not line:
                        # 空行表示段落结束
                        if current_paragraph:
                            paragraph_counter += 1
                            para_text = ' '.join(current_paragraph)
                            # 检查是否只是页码（过滤掉）
                            if not page_number_pattern.match(para_text):
                                paragraphs.append(para_text)
                                location_info.append((paragraph_counter, page_num, current_line_num))
                            current_paragraph = []
                            current_line_num = None
                        continue
                    
                    # 检查是否是单独的页码行（过滤掉）
                    if page_number_pattern.match(line):
                        continue
                    
                    # 提取行首的视觉行号
                    visual_line_num = None
                    match = line_number_pattern.match(line)
                    if match:
                        try:
                            visual_line_num = int(match.group(1).strip())
                            # 移除行首的数字
                            line = line[match.end():].lstrip()
                        except ValueError:
                            pass
                    
                    # 记录段落的第一个行号
                    if not current_paragraph and visual_line_num:
                        current_line_num = visual_line_num
                    
                    # 检查是否是新段落
                    is_new_para = False
                    if current_paragraph:
                        stripped = line.lstrip()
                        indent = len(line) - len(stripped)
                        if indent >= 4:
                            is_new_para = True
                        prev_line = current_paragraph[-1]
                        if prev_line and prev_line[-1] in '.。!?':
                            if stripped and not stripped[0].islower():
                                is_new_para = True
                    
                    if is_new_para and current_paragraph:
                        paragraph_counter += 1
                        para_text = ' '.join(current_paragraph)
                        if not page_number_pattern.match(para_text):
                            paragraphs.append(para_text)
                            location_info.append((paragraph_counter, page_num, current_line_num))
                        current_paragraph = [line]
                        current_line_num = visual_line_num
                    else:
                        current_paragraph.append(line)
                
                # 处理最后一个段落
                if current_paragraph:
                    paragraph_counter += 1
                    para_text = ' '.join(current_paragraph)
                    if not page_number_pattern.match(para_text):
                        paragraphs.append(para_text)
                        location_info.append((paragraph_counter, page_num, current_line_num))
            else:
                # 原始模式：每行单独处理
                for line in page_lines:
                    line = line.rstrip()
                    if not line or page_number_pattern.match(line):
                        continue
                    
                    # 提取行首的视觉行号
                    visual_line_num = None
                    match = line_number_pattern.match(line)
                    if match:
                        try:
                            visual_line_num = int(match.group(1).strip())
                            line = line[match.end():].lstrip()
                        except ValueError:
                            pass
                    
                    if line:
                        paragraph_counter += 1
                        paragraphs.append(line)
                        location_info.append((paragraph_counter, page_num, visual_line_num))
    
    return paragraphs, location_info


@register_reader('.pptx')
def read_pptx(path):
    """读取PPTX，返回内容列表和位置信息(段落号, 幻灯片号, None)"""
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    location_info = []
    paragraph_counter = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for ln in shape.text.splitlines():
                    line = ln.rstrip()
                    if line:  # 只处理非空行
                        paragraph_counter += 1
                        lines.append(line)
                        location_info.append((paragraph_counter, slide_num, None))
    
    return lines, location_info


def read_document(path, use_precise=True, merge_lines=True):
    """
    读取文档内容
    use_precise: 对docx文件是否使用高精度视觉行号计算
    merge_lines: 对pdf/txt是否合并连续行（将物理换行合并为逻辑段落）
    """
    ext = Path(path).suffix.lower()
    if ext not in READERS:
        raise ValueError(f"不支持的文件格式: {ext}。当前支持: {', '.join(READERS.keys())}")
    
    reader = READERS[ext]
    # 对不同格式传递相应参数
    if ext == '.docx':
        return reader(path, use_precise=use_precise)
    elif ext in ('.pdf', '.txt'):
        return reader(path, merge_lines=merge_lines)
    else:
        return reader(path)


def set_run_font(run, font_name='Times New Roman', east_asia='宋体', size=Pt(10), bold=False):
    """统一设置中英文字体"""
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn('w:eastAsia'), east_asia)
    run.font.size = size
    if bold:
        run.font.bold = True


def add_colored_run(paragraph, text, rgb, bold=False):
    run = paragraph.add_run(text)
    set_run_font(run, bold=bold)
    run.font.color.rgb = RGBColor(*rgb)
    return run


def add_page_number(paragraph):
    """在段落中添加页码字段"""
    run = paragraph.add_run()
    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')

    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = "PAGE"

    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'end')

    run._r.append(fldChar1)
    run._r.append(instrText)
    run._r.append(fldChar2)
    return run


def tokenize_text(text):
    """
    将文本拆分为单词和分隔符（空格/标点等）。
    保留所有内容以便后续重组。
    """
    return re.findall(r'\S+|\s+', text)


def split_sentences(text):
    """
    将文本按句子分割，保留分隔符。
    句子分隔符：中文句号、英文句号、问号、感叹号
    """
    if not text.strip():
        return []
    # 按句子分隔符分割，保留分隔符
    parts = re.split(r'([。！？.!?])', text)
    sentences = []
    current = ''
    for part in parts:
        current += part
        if re.match(r'[。！？.!?]$', part):
            # 遇到句子结束符，完成一个句子
            sentences.append(current)
            current = ''
    if current.strip():
        # 剩余内容（可能是不完整的句子）也作为一个句子
        sentences.append(current)
    return sentences


def word_diff_runs(text1, text2):
    """
    单词级精细差异分析，支持句子级缺失检测。
    返回 left_runs 和 right_runs，每项为 (text, is_diff, is_placeholder) 列表。
    is_placeholder=True 表示这是缺失内容占位符 [此处缺失句子]
    
    判断逻辑：
    1. 先将文本按句子分割
    2. 在句子级别对比
    3. 当检测到一边是完整句子（有分隔符），另一边没有对应句子时，显示占位符
    """
    # 首先尝试句子级对比
    sents1 = split_sentences(text1)
    sents2 = split_sentences(text2)
    
    # 如果都能分割成句子，进行句子级对比
    if len(sents1) > 0 and len(sents2) > 0:
        sm_sents = difflib.SequenceMatcher(None, sents1, sents2, autojunk=False)
        opcodes = sm_sents.get_opcodes()
        
        # 检查是否有句子级别的 delete/insert
        has_sentence_level_diff = any(
            tag in ('delete', 'insert') for tag, _, _, _, _ in opcodes
        )
        
        if has_sentence_level_diff:
            left_runs = []
            right_runs = []
            
            for tag, i1, i2, j1, j2 in opcodes:
                if tag == 'equal':
                    for k in range(i1, min(i2, len(sents1))):
                        left_runs.append((sents1[k], False, False))
                    for k in range(j1, min(j2, len(sents2))):
                        right_runs.append((sents2[k], False, False))
                elif tag == 'replace':
                    # 替换：两边内容不同
                    max_len = max(i2 - i1, j2 - j1)
                    for k in range(max_len):
                        if i1 + k < i2 and i1 + k < len(sents1):
                            left_runs.append((sents1[i1 + k], True, False))
                        if j1 + k < j2 and j1 + k < len(sents2):
                            right_runs.append((sents2[j1 + k], True, False))
                elif tag == 'delete':
                    # 左边有句子，右边缺失整句
                    for k in range(i1, min(i2, len(sents1))):
                        left_runs.append((sents1[k], True, False))
                        right_runs.append(('[此处缺失句子]', True, True))
                elif tag == 'insert':
                    # 右边有句子，左边缺失整句
                    for k in range(j1, min(j2, len(sents2))):
                        left_runs.append(('[此处缺失句子]', True, True))
                        right_runs.append((sents2[k], True, False))
            
            # 合并连续的 [此处缺失句子] 占位符
            def merge_consecutive_placeholders(runs):
                if not runs:
                    return runs
                merged = [runs[0]]
                for i in range(1, len(runs)):
                    prev_text, prev_diff, prev_placeholder = merged[-1]
                    curr_text, curr_diff, curr_placeholder = runs[i]
                    # 如果前一个是占位符且当前也是占位符，跳过当前（合并）
                    if prev_placeholder and curr_placeholder:
                        continue
                    merged.append((curr_text, curr_diff, curr_placeholder))
                return merged

            left_runs = merge_consecutive_placeholders(left_runs)
            right_runs = merge_consecutive_placeholders(right_runs)
            
            return left_runs, right_runs
    
    # 回退到单词级对比（不显示占位符）
    tokens1 = tokenize_text(text1)
    tokens2 = tokenize_text(text2)
    sm = difflib.SequenceMatcher(None, tokens1, tokens2, autojunk=False)

    left_runs = []
    right_runs = []

    for tag, i1, i2, j1, j2 in sm.get_opcodes():
        if tag == 'equal':
            seg1 = ''.join(tokens1[i1:i2])
            left_runs.append((seg1, False, False))
            seg2 = ''.join(tokens2[j1:j2])
            right_runs.append((seg2, False, False))
        elif tag == 'replace':
            seg1 = ''.join(tokens1[i1:i2])
            left_runs.append((seg1, True, False))
            seg2 = ''.join(tokens2[j1:j2])
            right_runs.append((seg2, True, False))
        elif tag == 'delete':
            seg1 = ''.join(tokens1[i1:i2])
            left_runs.append((seg1, True, False))
        elif tag == 'insert':
            seg2 = ''.join(tokens2[j1:j2])
            right_runs.append((seg2, True, False))

    if not left_runs:
        left_runs.append(('', False, False))
    if not right_runs:
        right_runs.append(('', False, False))

    return left_runs, right_runs


def get_word_line_number_offset(doc_path):
    """
    尝试读取 Word 文档的行号设置，返回行号起始偏移量。
    如果文档启用了行号，返回行号起始值；否则返回 None。
    """
    try:
        doc = Document(doc_path)
        for section in doc.sections:
            # 尝试获取行号设置
            sectPr = section._sectPr
            if hasattr(sectPr, 'lnNumType') and sectPr.lnNumType is not None:
                # 文档启用了行号
                lnNumType = sectPr.lnNumType
                start = getattr(lnNumType, 'start', 1)
                return int(start) if start else 1
    except Exception:
        pass
    return None


def build_diff_report(lines1, lines2, location_info1=None, location_info2=None):
    """
    使用 difflib.SequenceMatcher 分析差异，返回仅包含差异行的数据。
    每个元素为 (tag, left_loc, left_text, right_loc, right_text)
    tag 取值: replace, delete, insert
    location 格式: (paragraph_number, page_number) 或 None
    """
    sm = difflib.SequenceMatcher(None, lines1, lines2, autojunk=False)
    opcodes = sm.get_opcodes()
    
    # 如果没有提供位置信息，使用默认的段落索引+1
    if location_info1 is None:
        location_info1 = [(i + 1, 1) for i in range(len(lines1))]
    if location_info2 is None:
        location_info2 = [(i + 1, 1) for i in range(len(lines2))]

    rows = []
    for tag, i1, i2, j1, j2 in opcodes:
        if tag == 'equal':
            continue
        elif tag == 'replace':
            max_len = max(i2 - i1, j2 - j1)
            for k in range(max_len):
                left_exists = i1 + k < i2
                right_exists = j1 + k < j2
                left_loc = location_info1[i1 + k] if left_exists else None
                right_loc = location_info2[j1 + k] if right_exists else None
                ltext = lines1[i1 + k] if left_exists else ""
                rtext = lines2[j1 + k] if right_exists else ""
                if left_exists and right_exists:
                    rows.append(('replace', left_loc, ltext, right_loc, rtext))
                elif left_exists and not right_exists:
                    rows.append(('delete', left_loc, ltext, None, ""))
                elif not left_exists and right_exists:
                    rows.append(('insert', None, "", right_loc, rtext))
        elif tag == 'delete':
            for k in range(i2 - i1):
                left_loc = location_info1[i1 + k]
                ltext = lines1[i1 + k]
                rows.append(('delete', left_loc, ltext, None, ""))
        elif tag == 'insert':
            for k in range(j2 - j1):
                right_loc = location_info2[j1 + k]
                rtext = lines2[j1 + k]
                rows.append(('insert', None, "", right_loc, rtext))

    return rows


def set_cell_width(cell, width_inches):
    """通过底层 XML 强制设置单元格宽度"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcW = OxmlElement('w:tcW')
    tcW.set(qn('w:w'), str(int(width_inches * 1440)))
    tcW.set(qn('w:type'), 'dxa')
    tcPr.append(tcW)


def set_table_column_widths(table, widths_inches):
    """
    通过底层 XML 精确设置表格各列宽度及总宽度。
    widths_inches: 每列宽度的英寸值列表
    """
    tbl = table._tbl
    tblGrid = tbl.tblGrid
    gridCols = tblGrid.gridCol_lst
    for idx, w in enumerate(widths_inches):
        if idx < len(gridCols):
            gridCols[idx].set(qn('w:w'), str(int(w * 1440)))
            gridCols[idx].set(qn('w:type'), 'dxa')
        if idx < len(table.columns):
            table.columns[idx].width = Inches(w)

    tblPr = tbl.tblPr
    tblW_list = tblPr.xpath('./w:tblW')
    if tblW_list:
        tblW = tblW_list[0]
    else:
        tblW = OxmlElement('w:tblW')
        tblPr.append(tblW)
    total_width = sum(widths_inches)
    tblW.set(qn('w:w'), str(int(total_width * 1440)))
    tblW.set(qn('w:type'), 'dxa')


def generate_docx(rows, name1, name2, output_path):
    doc = Document()

    # 设置为横向（Landscape），并应用标准 1 英寸页边距
    section = doc.sections[0]
    section.orientation = WD_ORIENT.LANDSCAPE
    new_width, new_height = section.page_height, section.page_width
    section.page_width = new_width
    section.page_height = new_height
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)

    # 添加页脚页码
    footer = section.footer
    footer_para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = add_page_number(footer_para)
    set_run_font(run, size=Pt(9))

    # 颜色定义
    color_left = (0, 112, 192)    # 蓝色
    color_right = (255, 0, 0)     # 红色
    color_mark_replace = (255, 140, 0)  # 橙色
    color_gray = (100, 100, 100)
    color_placeholder = (0, 176, 80)  # 绿色，用于缺失句子

    # 第一行：文档对比报告
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_title.add_run('文档对比报告')
    set_run_font(run, east_asia='黑体', size=Pt(18), bold=False)

    # 第二行：两个文档名称
    p_names = doc.add_paragraph()
    p_names.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_names.add_run(name1)
    set_run_font(run, size=Pt(14), bold=True)
    run.font.color.rgb = RGBColor(*color_left)
    run = p_names.add_run(' VS ')
    set_run_font(run, size=Pt(14), bold=True)
    run = p_names.add_run(name2)
    set_run_font(run, size=Pt(14), bold=True)
    run.font.color.rgb = RGBColor(*color_right)

    # 第三行：生成日期时间
    now_str = datetime.now().strftime('%Y-%m-%d %H:%M')
    p_time = doc.add_paragraph()
    p_time.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_time.add_run(now_str)
    set_run_font(run, size=Pt(12), bold=True)
    run.font.color.rgb = RGBColor(*color_gray)

    # 时间行与图例行之间空一行
    doc.add_paragraph()

    # 图例说明（宋体）
    p = doc.add_paragraph()
    run = p.add_run('说明：')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('1. ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('= ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_left)
    run = p.add_run('为完全相同内容（已隐藏）  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('≠ ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_mark_replace)
    run = p.add_run('为有修改内容  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('- ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_left)
    run = p.add_run('为有删除内容  ')
    set_run_font(run, east_asia='宋体', bold=False)
    run = p.add_run('+ ')
    set_run_font(run, east_asia='宋体', bold=False)
    run.font.color.rgb = RGBColor(*color_right)
    run = p.add_run('为有新增内容')
    set_run_font(run, east_asia='宋体', bold=False)

    if not rows:
        p = doc.add_paragraph()
        run = p.add_run('两篇文档内容完全一致，未发现差异。')
        set_run_font(run, size=Pt(12))
        run.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.save(output_path)
        print(f"对比报告已保存至: {output_path}")
        return

    # 表格：5列（总宽度 8.8 英寸，适应 1 英寸边距的横向页面）
    table = doc.add_table(rows=1, cols=5)
    table.style = 'Table Grid'
    table.autofit = False
    table.allow_autofit = False

    col_widths = [0.55, 3.80, 0.45, 0.55, 3.80]
    set_table_column_widths(table, col_widths)

    # 设置表头 - 改为显示页码-段落号
    hdr_cells = table.rows[0].cells
    headers = ['位置', name1, '标记', '位置', name2]
    header_colors = [None, color_left, None, None, color_right]
    for idx, text in enumerate(headers):
        cell = hdr_cells[idx]
        set_cell_width(cell, col_widths[idx])
        p = cell.paragraphs[0]
        p.clear()
        run = p.add_run(text)
        set_run_font(run, bold=True)
        if header_colors[idx]:
            run.font.color.rgb = RGBColor(*header_colors[idx])
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for tag, left_loc, ltext, right_loc, rtext in rows:
        row_cells = table.add_row().cells

        # 位置1 - 格式: Page页码-L行号（或 Page页码-段落号）
        cell = row_cells[0]
        set_cell_width(cell, col_widths[0])
        p = cell.paragraphs[0]
        p.clear()
        if left_loc is not None:
            # 处理三元组 (para_num, page_num, line_num) 或二元组 (para_num, page_num)
            if len(left_loc) >= 3:
                para_num, page_num, line_num = left_loc[0], left_loc[1], left_loc[2]
                if line_num:
                    loc_text = f"Page{page_num}-L{line_num}"
                else:
                    loc_text = f"Page{page_num}-{para_num}"
            else:
                para_num, page_num = left_loc[0], left_loc[1]
                loc_text = f"Page{page_num}-{para_num}"
        else:
            loc_text = ""
        run = p.add_run(loc_text)
        set_run_font(run, size=Pt(8))
        run.font.color.rgb = RGBColor(*color_gray)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 位置2 - 格式: Page页码-L行号（或 Page页码-段落号）
        cell = row_cells[3]
        set_cell_width(cell, col_widths[3])
        p = cell.paragraphs[0]
        p.clear()
        if right_loc is not None:
            # 处理三元组或二元组
            if len(right_loc) >= 3:
                para_num, page_num, line_num = right_loc[0], right_loc[1], right_loc[2]
                if line_num:
                    loc_text = f"Page{page_num}-L{line_num}"
                else:
                    loc_text = f"Page{page_num}-{para_num}"
            else:
                para_num, page_num = right_loc[0], right_loc[1]
                loc_text = f"Page{page_num}-{para_num}"
        else:
            loc_text = ""
        run = p.add_run(loc_text)
        set_run_font(run, size=Pt(8))
        run.font.color.rgb = RGBColor(*color_gray)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 内容1
        cell = row_cells[1]
        set_cell_width(cell, col_widths[1])
        p1 = cell.paragraphs[0]
        p1.clear()
        if tag == 'replace':
            runs1, _ = word_diff_runs(ltext, rtext)
            for text_seg, is_diff, is_placeholder in runs1:
                if is_placeholder:
                    # 缺失占位符：绿色粗体
                    add_colored_run(p1, text_seg, color_placeholder, bold=True)
                elif is_diff:
                    add_colored_run(p1, text_seg, color_left, bold=True)
                else:
                    run = p1.add_run(text_seg)
                    set_run_font(run)
        elif tag == 'delete':
            add_colored_run(p1, ltext, color_left, bold=True)
            # 右侧整行空白，不添加任何标识
        elif tag == 'insert':
            # 左侧整行空白，不添加任何标识
            pass
        else:
            run = p1.add_run(ltext)
            set_run_font(run)

        # 标记
        cell = row_cells[2]
        set_cell_width(cell, col_widths[2])
        p_mark = cell.paragraphs[0]
        p_mark.clear()
        p_mark.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if tag == 'replace':
            add_colored_run(p_mark, '≠', color_mark_replace, bold=True)
        elif tag == 'delete':
            add_colored_run(p_mark, '-', color_left, bold=True)
        elif tag == 'insert':
            add_colored_run(p_mark, '+', color_right, bold=True)

        # 内容2
        cell = row_cells[4]
        set_cell_width(cell, col_widths[4])
        p2 = cell.paragraphs[0]
        p2.clear()
        if tag == 'replace':
            _, runs2 = word_diff_runs(ltext, rtext)
            for text_seg, is_diff, is_placeholder in runs2:
                if is_placeholder:
                    # 缺失占位符：绿色粗体
                    add_colored_run(p2, text_seg, color_placeholder, bold=True)
                elif is_diff:
                    add_colored_run(p2, text_seg, color_right, bold=True)
                else:
                    run = p2.add_run(text_seg)
                    set_run_font(run)
        elif tag == 'insert':
            add_colored_run(p2, rtext, color_right, bold=True)
        elif tag == 'delete':
            # 删除行：右侧整行空白，不添加任何标识
            pass
        else:
            run = p2.add_run(rtext)
            set_run_font(run)

    doc.save(output_path)
    print(f"对比报告已保存至: {output_path}")


def main():
    import argparse
    
    parser = argparse.ArgumentParser(description='文档对比工具')
    parser.add_argument('file1', help='第一个文档路径')
    parser.add_argument('file2', help='第二个文档路径')
    parser.add_argument('--calibrate', action='store_true', help='校准模式：输出段落位置信息用于调试')
    parser.add_argument('--no-merge', action='store_true', help='PDF/TXT文件不合并连续行（按原始行对比）')
    
    args = parser.parse_args()
    
    file1 = args.file1
    file2 = args.file2
    use_precise = True
    merge_lines = not args.no_merge  # 默认合并，--no-merge 时关闭
    
    if not os.path.exists(file1):
        print(f"错误: 文件不存在: {file1}")
        sys.exit(1)
    if not os.path.exists(file2):
        print(f"错误: 文件不存在: {file2}")
        sys.exit(1)

    name1 = Path(file1).stem
    name2 = Path(file2).stem
    output_name = f"Comparison_{name1}_VS_{name2}.docx"
    output_path = os.path.join(os.getcwd(), output_name)

    print(f"正在读取 {file1} ...")
    result1 = read_document(file1, use_precise=use_precise, merge_lines=merge_lines)
    if isinstance(result1, tuple):
        lines1, location_info1 = result1
    else:
        lines1 = result1
        location_info1 = None
    print(f"  共 {len(lines1)} 段落")
    
    # 校准模式：显示前几个段落的位置信息
    if args.calibrate and location_info1:
        print("\n校准信息（前5个段落）:")
        for i in range(min(5, len(lines1))):
            loc = location_info1[i]
            para_num, page_num, line_num = loc[0], loc[1], loc[2] if len(loc) > 2 else None
            if line_num:
                print(f"  段落 {i+1}: Page{page_num}-L{line_num}")
            else:
                print(f"  段落 {i+1}: Page{page_num}-{para_num}")
        print()

    print(f"正在读取 {file2} ...")
    result2 = read_document(file2, use_precise=use_precise, merge_lines=merge_lines)
    if isinstance(result2, tuple):
        lines2, location_info2 = result2
    else:
        lines2 = result2
        location_info2 = None
    print(f"  共 {len(lines2)} 段落")

    print("正在分析差异 ...")
    rows = build_diff_report(lines1, lines2, location_info1, location_info2)

    print("正在生成报告 ...")
    generate_docx(rows, name1, name2, output_path)

    diff_count = len(rows)
    print(f"差异行数: {diff_count}")
    print(f"对比报告已保存: {output_path}")


if __name__ == '__main__':
    main()
